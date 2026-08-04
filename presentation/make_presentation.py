"""Generates KeyForge_Presentation.pptx — full project presentation deck."""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

_ANCHORS = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT = HERE / "KeyForge_Presentation.pptx"
DB_DIAGRAM = ROOT / "Diagrams  DataBase" / "keyforge_28_07_2026.png"

# ---- Brand palette ----
INK = RGBColor(0x16, 0x18, 0x1D)
INK_SOFT = RGBColor(0x1E, 0x22, 0x29)
SUBTLE = RGBColor(0x6E, 0x77, 0x83)
MUTED = RGBColor(0x9A, 0xA3, 0xAD)
ACCENT = RGBColor(0x3D, 0x94, 0xFF)
ACCENT_DARK = RGBColor(0x00, 0x7B, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE5, 0xE5)
SUCCESS = RGBColor(0x3E, 0xCF, 0x78)
WARNING = RGBColor(0xDB, 0xA6, 0x37)
DANGER = RGBColor(0xE0, 0x58, 0x49)
GREEN_DEEP = RGBColor(0x24, 0x54, 0x3C)

FONT = "Segoe UI"
FONT_BRAND = "Georgia"

SW, SH = Inches(13.333), Inches(7.5)


def _set_spacing(run, value):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(value))


def add_text(slide, x, y, w, h, runs, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor="top", spacing=None, line_spacing=None,
             font=FONT, gap=None):
    """runs: str, or list of (text, dict-overrides), or list of paragraphs (list of runs)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = _ANCHORS.get(anchor, MSO_ANCHOR.TOP)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    elif isinstance(runs, tuple):
        runs = [[runs]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if gap:
            p.space_before = Pt(gap)
        if line_spacing:
            p.line_spacing = line_spacing
        for text, ov in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.color.rgb = ov.get("color", color)
            f.italic = ov.get("italic", False)
            if ov.get("spacing"):
                _set_spacing(r, ov["spacing"])
    return box


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def add_pill(slide, x, y, w, h, text, fill, fg=CARD, size=12, bold=True):
    sp = add_rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg
    return sp


def add_badge(slide, x, y, w, h, text, size=11):
    return add_pill(slide, x, y, w, h, text, ACCENT, size=size)


def slide_header(slide, kicker, title, accent=ACCENT):
    add_text(slide, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.3),
             kicker, size=11, color=accent, bold=True, spacing=300)
    add_text(slide, Inches(0.7), Inches(0.7), Inches(11.9), Inches(0.75),
             title, size=28, color=INK, bold=True)
    add_rect(slide, Inches(0.72), Inches(1.42), Inches(0.95), Inches(0.055), fill=accent)
    add_text(slide, Inches(0.7), Inches(1.58), Inches(11.9), Inches(0.3),
             "", size=6)


def slide_footer(slide, index, total):
    add_rect(slide, Inches(0.7), Inches(7.12), Inches(11.93), Pt(1), fill=BORDER)
    add_text(slide, Inches(0.7), Inches(7.2), Inches(6), Inches(0.25),
             "KeyForge · Final Project · Web Development", size=9, color=MUTED)
    add_text(slide, Inches(11.6), Inches(7.2), Inches(1.03), Inches(0.25),
             f"{index:02d} / {total:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, size=13.5, gap=6, color=INK, marker_color=ACCENT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (head, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        r = p.add_run()
        r.text = "▸ "
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = marker_color
        if head:
            r2 = p.add_run()
            r2.text = head
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.bold = True
            r2.font.color.rgb = color
        if rest:
            r3 = p.add_run()
            r3.text = rest
            r3.font.name = FONT
            r3.font.size = Pt(size)
            r3.font.color.rgb = color
    return box


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def blank_slide(prs, bg=CARD):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, fill=bg)
    return s


# --------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
TOTAL = 20

# ---- 1. Title -------------------------------------------------------------
s = blank_slide(prs, bg=INK)
add_rect(s, 0, Inches(4.35), SW, Pt(1), fill=RGBColor(0x2B, 0x31, 0x3A))
add_rect(s, Inches(0.85), Inches(2.18), Inches(1.1), Inches(0.07), fill=ACCENT)
add_text(s, Inches(0.85), Inches(2.38), Inches(11.6), Inches(1.2),
         "KeyForge", size=64, color=CARD, bold=True, font=FONT_BRAND)
add_text(s, Inches(0.88), Inches(3.5), Inches(11.6), Inches(0.6),
         "Mechanical Keyboard E-Commerce Platform", size=22, color=MUTED)
add_text(s, Inches(0.88), Inches(4.6), Inches(11.6), Inches(1.6),
         [
             [("A Full-Stack Web Application — React · FastAPI · MongoDB", {"size": 15, "color": MUTED})],
             [("Final Project Presentation", {"size": 15, "color": MUTED})],
         ], gap=6)
add_text(s, Inches(0.88), Inches(6.35), Inches(11.6), Inches(0.8),
         [
             [("Presented by:  ", {"bold": True, "color": CARD}),
              ("Kimhong Pech", {"color": CARD})],
             [("Live:  ", {"bold": True, "color": CARD}),
              ("https://key-forge-final-project.vercel.app", {"color": ACCENT}),
              ("     Repo:  ", {"bold": True, "color": CARD}),
              ("github.com/kimhongpech-dev/KeyForge-Final-Project", {"color": ACCENT})],
         ], size=12.5)

# ---- 2. Agenda ------------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "OUTLINE", "Agenda")
items = [
    ("01", "Project Overview & Goals", "What KeyForge is, who it serves, and what it achieves"),
    ("02", "Problem & Motivation", "Why build a full e-commerce platform from scratch"),
    ("03", "Technology Stack", "Frontend, backend, database, and deployment choices"),
    ("04", "System Architecture", "How the pieces talk to each other"),
    ("05", "Database & API Design", "MongoDB data model and REST endpoints"),
    ("06", "Core Features", "Catalog, auth, cart & checkout, orders, admin dashboard"),
    ("07", "UX & Design", "Theming, responsiveness, accessibility"),
    ("08", "Deployment & Challenges", "GitHub, Vercel, MongoDB Atlas, problems solved"),
    ("09", "Testing, Demo & Future Work", "How it was verified, demo flow, next steps"),
]
y = Inches(1.85)
for num, head, sub in items:
    add_text(s, Inches(0.7), y, Inches(0.7), Inches(0.4), num, size=18, color=ACCENT, bold=True)
    add_text(s, Inches(1.45), y - Inches(0.03), Inches(7.2), Inches(0.35), head, size=15, color=INK, bold=True)
    add_text(s, Inches(1.45), y + Inches(0.3), Inches(11.0), Inches(0.3), sub, size=11.5, color=SUBTLE)
    y += Inches(0.55)
slide_footer(s, 2, TOTAL)

# ---- 3. Project overview --------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 1 · VISION", "Project Overview")
add_text(s, Inches(0.7), Inches(1.85), Inches(7.4), Inches(1.1),
         [[("KeyForge is a production-style online store for mechanical keyboards, "
            "gaming mice, and earphones — built end-to-end, from product browsing to "
            "admin analytics.", {"size": 14, "color": INK})]], line_spacing=1.25)
add_text(s, Inches(0.7), Inches(2.75), Inches(7.4), Inches(0.3),
         "Key Goals", size=13, color=ACCENT_DARK, bold=True)
bullets(s, Inches(0.7), Inches(3.1), Inches(7.4), Inches(3.5), [
    ("Complete shopping journey: ", "browse → search → cart → checkout → track orders"),
    ("Real accounts & security: ", "registration, login, JWT sessions, admin roles"),
    ("Business-grade stock control: ", "inventory enforced atomically at checkout"),
    ("Full admin operations: ", "product CRUD, stock edits, order workflow, analytics"),
    ("Production deployment: ", "live on Vercel with MongoDB Atlas, auto-deploy on push"),
])
cards = [
    ("25", "Products seeded", "keyboards, mice, earphones"),
    ("3", "Categories", "with live search & filters"),
    ("4", "Order stages", "pending → confirmed → shipped → delivered"),
    ("6", "Charts & stats", "custom-built SVG analytics"),
]
x = Inches(8.35)
for num, head, sub in cards:
    add_rect(s, x, Inches(1.85), Inches(4.28), Inches(1.28), fill=LIGHT_BG, line=BORDER, line_w=0.75)
    add_text(s, x + Inches(0.25), Inches(2.0), Inches(3.8), Inches(0.5), num, size=24, color=ACCENT, bold=True)
    add_text(s, x + Inches(0.25), Inches(2.5), Inches(3.8), Inches(0.3), head, size=12.5, color=INK, bold=True)
    add_text(s, x + Inches(0.25), Inches(2.78), Inches(3.8), Inches(0.3), sub, size=10.5, color=SUBTLE)
    y += Inches(0)
    x = x
add_rect(s, Inches(8.35), Inches(3.28), Inches(4.28), Inches(1.28), fill=LIGHT_BG, line=BORDER, line_w=0.75)
add_text(s, Inches(8.6), Inches(3.43), Inches(3.8), Inches(0.5), "2", size=24, color=ACCENT, bold=True)
add_text(s, Inches(8.6), Inches(3.93), Inches(3.8), Inches(0.3), "Theme modes", size=12.5, color=INK, bold=True)
add_text(s, Inches(8.6), Inches(4.21), Inches(3.8), Inches(0.3), "light & dark, remembers choice", size=10.5, color=SUBTLE)
add_rect(s, Inches(8.35), Inches(4.71), Inches(4.28), Inches(1.28), fill=LIGHT_BG, line=BORDER, line_w=0.75)
add_text(s, Inches(8.6), Inches(4.86), Inches(3.8), Inches(0.5), "1", size=24, color=ACCENT, bold=True)
add_text(s, Inches(8.6), Inches(5.36), Inches(3.8), Inches(0.3), "Owner / admin account", size=12.5, color=INK, bold=True)
add_text(s, Inches(8.6), Inches(5.64), Inches(3.8), Inches(0.3), "boss@keyforge.com · promoted via script", size=10.5, color=SUBTLE)
add_rect(s, Inches(8.35), Inches(6.14), Inches(4.28), Inches(0.85), fill=RGBColor(0xEA, 0xF2, 0xFB))
add_text(s, Inches(8.6), Inches(6.28), Inches(3.8), Inches(0.6),
         [[("Deployed live at ", {"color": INK_SOFT}),
           ("key-forge-final-project.vercel.app", {"color": ACCENT_DARK, "bold": True})]],
         size=11.5)
slide_footer(s, 3, TOTAL)

# ---- 4. Problem & motivation ---------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 1 · VISION", "Problem & Motivation")
add_text(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.9),
         [[("A shopping site is a small operating system: products, money, identity, and inventory "
            "must all stay consistent in real time. Most course projects stop at mock data — "
            "KeyForge goes the full distance to a working, deployed product.", {"size": 14})]],
         line_spacing=1.25)
cols = [
    ("The Problem", ACCENT_DARK, [
        ("Fake checkouts", "orders that never reduce stock or create records"),
        ("No real accounts", "everyone sees everything, no roles or ownership"),
        ("Manual operations", "no way to manage products, stock, or orders"),
        ("Static catalog", "products hard-coded in the frontend, unchangeable"),
        ("No deployment", "works on one laptop only, impossible to share"),
    ]),
    ("How KeyForge Solves It", SUCCESS, [
        ("Real checkout", "orders persisted, stock decremented atomically"),
        ("JWT accounts & roles", "customers vs. admin, protected APIs & pages"),
        ("Admin dashboard", "full product CRUD, stock, order workflow, analytics"),
        ("Database-driven catalog", "25 products served from MongoDB, editable"),
        ("Production deployment", "live on Vercel, MongoDB Atlas, CI via git push"),
    ]),
]
x = Inches(0.7)
for title, accent, rows in cols:
    add_rect(s, x, Inches(2.9), Inches(5.85), Inches(3.9), fill=LIGHT_BG, line=BORDER)
    add_rect(s, x, Inches(2.9), Inches(5.85), Inches(0.55), fill=accent)
    add_text(s, x + Inches(0.3), Inches(2.99), Inches(5.2), Inches(0.4), title, size=15, color=CARD, bold=True)
    yy = Inches(3.65)
    for head, sub in rows:
        add_text(s, x + Inches(0.3), yy, Inches(5.3), Inches(0.3), head, size=12, color=INK, bold=True)
        add_text(s, x + Inches(0.3), yy + Inches(0.25), Inches(5.3), Inches(0.3), sub, size=10.5, color=SUBTLE)
        yy += Inches(0.6)
    x += Inches(6.05)
slide_footer(s, 4, TOTAL)

# ---- 5. Tech stack --------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "Technology Stack")
groups = [
    ("Frontend", ACCENT, [
        ("React 19 + Vite", "fast builds, HMR, component model"),
        ("React Router", "7 pages, dynamic product routes"),
        ("Context API", "auth, cart & theme state — no extra libraries"),
        ("Custom CSS", "design-system variables, responsive grid"),
    ]),
    ("Backend", ACCENT_DARK, [
        ("Python 3.12 + FastAPI", "async, typed, auto-documented"),
        ("Motor (async driver)", "non-blocking MongoDB access"),
        ("PyJWT + bcrypt", "signed sessions, hashed passwords"),
        ("Pydantic schemas", "validated request/response models"),
    ]),
    ("Data & Ops", SUCCESS, [
        ("MongoDB Atlas", "cloud database, free tier, auto-indexed"),
        ("GitHub", "private repo, versioned history"),
        ("Vercel", "static site + Python serverless functions"),
        ("GitHub ↔ Vercel", "auto-deploy on every push to main"),
    ]),
]
x = Inches(0.7)
for title, accent, rows in groups:
    add_rect(s, x, Inches(1.9), Inches(3.83), Inches(4.7), fill=LIGHT_BG, line=BORDER)
    add_rect(s, x, Inches(1.9), Inches(3.83), Inches(0.62), fill=accent)
    add_text(s, x + Inches(0.25), Inches(2.0), Inches(3.3), Inches(0.45), title, size=16, color=CARD, bold=True)
    yy = Inches(2.75)
    for head, sub in rows:
        add_text(s, x + Inches(0.25), yy, Inches(3.35), Inches(0.3), head, size=12.5, color=INK, bold=True)
        add_text(s, x + Inches(0.25), yy + Inches(0.27), Inches(3.35), Inches(0.55), sub, size=10.5, color=SUBTLE)
        yy += Inches(0.88)
    x += Inches(4.0)
add_text(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.3),
         [[("No heavy UI libraries — charts, theming, and layout are hand-built for full control and "
            "learning value.", {"size": 11.5, "italic": True, "color": SUBTLE})]])
slide_footer(s, 5, TOTAL)

# ---- 6. Architecture ------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "System Architecture")
layer1 = [
    ("User / Browser", RGBColor(0x2B, 0x31, 0x3A), CARD),
]
add_rect(s, Inches(1.0), Inches(2.1), Inches(2.6), Inches(1.5), fill=INK)
add_text(s, Inches(1.0), Inches(2.75), Inches(2.6), Inches(0.5), "User / Browser", size=15, color=CARD, bold=True, align=PP_ALIGN.CENTER)
add_rect(s, Inches(3.75), Inches(2.55), Inches(0.9), Inches(0.6), fill=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
add_text(s, Inches(3.85), Inches(2.66), Inches(0.8), Inches(0.4), "HTTPS", size=9, color=CARD, bold=True, align=PP_ALIGN.CENTER)
add_rect(s, Inches(4.85), Inches(1.7), Inches(7.8), Inches(3.4), fill=LIGHT_BG, line=ACCENT, line_w=1.5)
add_text(s, Inches(5.05), Inches(1.8), Inches(7.4), Inches(0.35), "Vercel — one deployment, two runtimes", size=13, color=ACCENT_DARK, bold=True)
add_rect(s, Inches(5.05), Inches(2.2), Inches(3.5), Inches(1.55), fill=CARD, line=BORDER)
add_text(s, Inches(5.2), Inches(2.32), Inches(3.2), Inches(0.3), "React SPA (static)", size=12, color=INK, bold=True)
add_text(s, Inches(5.2), Inches(2.66), Inches(3.2), Inches(1.0),
         [[("Pages · cart · orders ·", {"size": 10.5, "color": SUBTLE})],
          [("admin dashboard · theme", {"size": 10.5, "color": SUBTLE})],
          [("served as static files", {"size": 10.5, "color": SUBTLE})]])
add_rect(s, Inches(8.85), Inches(2.2), Inches(3.6), Inches(1.55), fill=CARD, line=BORDER)
add_text(s, Inches(9.0), Inches(2.32), Inches(3.3), Inches(0.3), "FastAPI (Python functions)", size=12, color=INK, bold=True)
add_text(s, Inches(9.0), Inches(2.66), Inches(3.3), Inches(1.0),
         [[("REST API under /api — auth,", {"size": 10.5, "color": SUBTLE})],
          [("products, orders, admin", {"size": 10.5, "color": SUBTLE})],
          [("runs as serverless functions", {"size": 10.5, "color": SUBTLE})]])
add_text(s, Inches(5.05), Inches(4.0), Inches(7.4), Inches(0.5),
         "SPA calls the same domain — no CORS issues in production", size=10.5, color=SUBTLE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(8.85), Inches(4.45), Inches(0.9), Inches(0.6), fill=SUCCESS, shape=MSO_SHAPE.RIGHT_ARROW)
add_rect(s, Inches(5.05), Inches(4.45), Inches(0.9), Inches(0.6), fill=SUCCESS, shape=MSO_SHAPE.RIGHT_ARROW)
add_rect(s, Inches(5.05), Inches(5.2), Inches(7.6), Inches(1.15), fill=INK)
add_text(s, Inches(5.25), Inches(5.42), Inches(7.2), Inches(0.35), "MongoDB Atlas (cloud)", size=14, color=CARD, bold=True)
add_text(s, Inches(5.25), Inches(5.8), Inches(7.2), Inches(0.4),
         "users · products · orders  —  async driver, indexed lookups", size=11, color=MUTED)
add_text(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.4),
         [[("Key flows:  ", {"bold": True, "color": INK}),
           ("JWT issued at login and verified per request · cart lives in browser state · "
            "orders & stock updates are atomic database operations.", {"color": SUBTLE})]], size=11.5)
slide_footer(s, 6, TOTAL)

# ---- 7. Database design ---------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "Database Design — MongoDB Atlas")
add_text(s, Inches(0.7), Inches(1.78), Inches(11.9), Inches(0.35),
         [[("Three collections, one database. Order items are embedded (no join needed); "
            "stock lives on the product so it can be decremented atomically.", {"size": 12.5, "color": SUBTLE})]])
add_rect(s, Inches(0.7), Inches(2.25), Inches(3.9), Inches(2.1), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(0.95), Inches(2.4), Inches(3.4), Inches(0.3), "users", size=14, color=ACCENT_DARK, bold=True)
add_text(s, Inches(0.95), Inches(2.78), Inches(3.4), Inches(1.4),
         [[("name, email (unique)", {"size": 11})],
          [("password — bcrypt hash", {"size": 11})],
          [("is_admin flag", {"size": 11})],
          [("created_at", {"size": 11})]])
add_rect(s, Inches(4.8, ) if False else Inches(4.8), Inches(2.25), Inches(3.9), Inches(2.1), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(5.05), Inches(2.4), Inches(3.4), Inches(0.3), "products", size=14, color=ACCENT_DARK, bold=True)
add_text(s, Inches(5.05), Inches(2.78), Inches(3.4), Inches(1.4),
         [[("id, name, price", {"size": 11})],
          [("category, description, image", {"size": 11})],
          [("stock — source of truth", {"size": 11})],
          [("25 seeded records", {"size": 11})]])
add_rect(s, Inches(8.9), Inches(2.25), Inches(3.7), Inches(2.1), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(9.15), Inches(2.4), Inches(3.2), Inches(0.3), "orders", size=14, color=ACCENT_DARK, bold=True)
add_text(s, Inches(9.15), Inches(2.78), Inches(3.2), Inches(1.4),
         [[("user email + items embedded", {"size": 11})],
          [("total, status, timestamps", {"size": 11})],
          [("tracking: 4 status values", {"size": 11})]])
add_text(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.3),
         "Entity-relationship overview (from the project's design phase):", size=12.5, color=INK, bold=True)
if DB_DIAGRAM.exists():
    with Image.open(DB_DIAGRAM) as im:
        iw, ih = im.size
    pic_w = Inches(11.9)
    ratio = ih / iw
    pic_h = Emu(int(pic_w * ratio))
    if pic_h > Inches(2.1):
        pic_h = Inches(2.1)
        pic_w = Emu(int(pic_h / ratio))
    s.shapes.add_picture(str(DB_DIAGRAM), Inches(0.7), Inches(4.95), width=pic_w, height=pic_h)
else:
    add_text(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.4), "diagram image not found", size=11, color=DANGER)
slide_footer(s, 7, TOTAL)

# ---- 8. API design --------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "REST API Design")
rows = [
    ("GET", "/api/products", "List products · optional ?search & ?category filters"),
    ("GET", "/api/products/{id}", "Single product detail"),
    ("GET", "/api/products/categories", "Distinct category list for filters"),
    ("POST", "/api/auth/register", "Create account (email unique, password hashed)"),
    ("POST", "/api/auth/login", "Verify credentials, return JWT token"),
    ("GET", "/api/auth/me", "Resolve current user from token"),
    ("POST", "/api/orders", "Place order · atomic stock decrement · rejects short stock"),
    ("GET", "/api/orders", "List my orders (newest first)"),
    ("POST", "/api/orders/{id}/cancel", "Client cancellation · restores stock"),
    ("GET", "/api/admin/products", "Admin product list"),
    ("POST / PUT / DELETE", "/api/admin/products/{id}", "Admin product CRUD & stock updates"),
    ("GET / PUT", "/api/admin/orders/{id}", "Admin order list & status workflow"),
    ("GET", "/api/admin/stats", "Dashboard analytics (revenue, statuses, categories, top)"),
]
tbl_shape = s.shapes.add_table(len(rows) + 1, 3, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.6))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.1)
tbl.columns[1].width = Inches(4.1)
tbl.columns[2].width = Inches(5.7)
hdr = ["Method", "Endpoint", "Purpose"]
for c, htext in enumerate(hdr):
    cell = tbl.cell(0, c)
    cell.text = htext
    cell.fill.solid()
    cell.fill.fore_color.rgb = INK
    p = cell.text_frame.paragraphs[0]
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = CARD
    cell.margin_left = Pt(8)
    cell.margin_top = Pt(4)
for i, (m, ep, desc) in enumerate(rows, start=1):
    vals = [m, ep, desc]
    for c, v in enumerate(vals):
        cell = tbl.cell(i, c)
        cell.text = v
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD if i % 2 else RGBColor(0xF7, 0xF8, 0xFA)
        p = cell.text_frame.paragraphs[0]
        r = p.runs[0]
        r.font.name = FONT
        r.font.size = Pt(10.5)
        r.font.color.rgb = INK
        if c == 0:
            r.font.bold = True
            r.font.color.rgb = ACCENT_DARK
        if c == 1:
            r.font.name = "Consolas"
            r.font.size = Pt(10)
        cell.margin_left = Pt(8)
        cell.margin_top = Pt(2)
add_text(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.4),
         [[("All errors return a consistent shape — ", {"size": 11.5, "color": SUBTLE}),
           ("""{"error": "<message>"}""", {"size": 11.5, "color": DANGER, "font": "Consolas"}),
           (" — handled by one global exception handler.", {"size": 11.5, "color": SUBTLE})]])
slide_footer(s, 8, TOTAL)

# ---- 9. Catalog -----------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "Feature — Product Catalog & Discovery")
left = [
    ("Live search", "matches product name & description on the server"),
    ("Category filters", "Keybord · Mouse · Earsphone — one click filtering"),
    ("Stock awareness", "cards show live stock; out-of-stock items can't be added"),
    ("Product details page", "full description, price, and add-to-cart with stock caps"),
]
x = Inches(0.7)
add_rect(s, x, Inches(1.9), Inches(5.9), Inches(4.3), fill=LIGHT_BG, line=BORDER)
add_text(s, x + Inches(0.3), Inches(2.1), Inches(5.3), Inches(0.35), "What the customer sees", size=13, color=ACCENT_DARK, bold=True)
yy = Inches(2.55)
for head, sub in left:
    add_text(s, x + Inches(0.3), yy, Inches(5.3), Inches(0.3), f"▸ {head}", size=12, color=INK, bold=True)
    add_text(s, x + Inches(0.55), yy + Inches(0.26), Inches(5.1), Inches(0.55), sub, size=11, color=SUBTLE)
    yy += Inches(0.78)
add_rect(s, Inches(6.85), Inches(1.9), Inches(5.78), Inches(4.3), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(7.15), Inches(2.1), Inches(5.2), Inches(0.35), "How it's implemented", size=13, color=ACCENT_DARK, bold=True)
bullets(s, Inches(7.15), Inches(2.55), Inches(5.2), Inches(3.5), [
    ("Server-side filtering: ", "FastAPI builds a MongoDB query from query params"),
    ("Cache in memory: ", "first fetch fills a module-level cache per session"),
    ("Reusable components: ", "ProductCard, SearchBar, CategoryFilter"),
    ("Clean state: ", "one products[] state driven by search / category effects"),
], size=11.5, gap=8)
add_text(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
         [[("Design detail: ", {"bold": True, "color": INK}),
           ("the API returns JSON with stock included, so every card always shows the "
            "current inventory — no stale numbers after an order.", {"color": SUBTLE})]], size=11.5)
slide_footer(s, 9, TOTAL)

# ---- 10. Auth -------------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "Feature — Authentication & Authorization")
flow = [
    ("Register / Login", "POST /api/auth/register | /login", ACCENT),
    ("JWT issued", "signed token with email + role claims", ACCENT_DARK),
    ("Client stores", "token kept in memory + localStorage", WARNING),
    ("Protected APIs", "verify bearer token on every request", SUCCESS),
]
x = Inches(0.7)
for head, sub, color in flow:
    add_rect(s, x, Inches(1.95), Inches(2.72), Inches(1.5), fill=LIGHT_BG, line=BORDER)
    add_rect(s, x, Inches(1.95), Inches(2.72), Inches(0.12), fill=color)
    add_text(s, x + Inches(0.18), Inches(2.2), Inches(2.36), Inches(0.55), head, size=12, color=INK, bold=True)
    add_text(s, x + Inches(0.18), Inches(2.72), Inches(2.36), Inches(0.6), sub, size=9.5, color=SUBTLE)
    if x + 2.72 < Inches(11):
        ar = add_rect(s, x + Inches(2.72), Inches(2.45), Inches(0.38), Inches(0.5), fill=ACCENT, shape=MSO_SHAPE.CHEVRON)
    x += Inches(3.1)
bullets(s, Inches(0.7), Inches(3.85), Inches(6.0), Inches(2.6), [
    ("Passwords never stored in plain text", "bcrypt hashing with per-user salt"),
    ("Stateless sessions", "PyJWT signs the token; no server session storage"),
    ("Admin flag in token", "the API trusts role claims, the UI hides admin links"),
    ("Edge cases handled", "duplicate email → 400, wrong password → 401, expired token → re-login"),
], size=12, gap=8)
add_rect(s, Inches(6.95), Inches(3.85), Inches(5.68), Inches(2.6), fill=INK)
add_text(s, Inches(7.25), Inches(4.05), Inches(5.1), Inches(0.35), "Role model", size=13, color=CARD, bold=True)
add_rect(s, Inches(7.25), Inches(4.5), Inches(2.45), Inches(1.55), fill=RGBColor(0x2B, 0x31, 0x3A))
add_text(s, Inches(7.4), Inches(4.62), Inches(2.15), Inches(0.3), "Customer", size=12, color=CARD, bold=True)
add_text(s, Inches(7.4), Inches(4.95), Inches(2.15), Inches(1.0),
         [[("browse · cart · order", {"size": 10, "color": MUTED})],
          [("track & cancel own orders", {"size": 10, "color": MUTED})]])
add_rect(s, Inches(9.9), Inches(4.5), Inches(2.45), Inches(1.55), fill=RGBColor(0x2B, 0x31, 0x3A))
add_text(s, Inches(10.05), Inches(4.62), Inches(2.15), Inches(0.3), "Admin", size=12, color=CARD, bold=True)
add_text(s, Inches(10.05), Inches(4.95), Inches(2.15), Inches(1.0),
         [[("everything a customer can", {"size": 10, "color": MUTED})],
          [("+ dashboard, CRUD, workflow", {"size": 10, "color": MUTED})]])
add_text(s, Inches(6.95), Inches(6.35), Inches(5.68), Inches(0.4),
         "Demo account: boss@keyforge.com (admin)", size=10.5, color=MUTED)
slide_footer(s, 10, TOTAL)

# ---- 11. Cart & checkout --------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "Feature — Cart & Checkout")
left = [
    ("Global cart state", "Context provider so every page shares one cart"),
    ("Quantity controls", "add, remove, and adjust — capped at available stock"),
    ("Live totals", "subtotal recalculated from cart items automatically"),
    ("Confirmation screen", "order number, summary, and items after success"),
]
add_rect(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(3.2), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(0.35), "Customer experience", size=13, color=ACCENT_DARK, bold=True)
yy = Inches(2.55)
for head, sub in left:
    add_text(s, Inches(1.0), yy, Inches(5.3), Inches(0.3), f"▸ {head}", size=12, color=INK, bold=True)
    add_text(s, Inches(1.25), yy + Inches(0.26), Inches(5.1), Inches(0.5), sub, size=11, color=SUBTLE)
    yy += Inches(0.72)
add_rect(s, Inches(6.85), Inches(1.9), Inches(5.78), Inches(3.2), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(7.15), Inches(2.1), Inches(5.2), Inches(0.35), "Stock safety — the interesting part", size=13, color=SUCCESS, bold=True)
add_text(s, Inches(7.15), Inches(2.55), Inches(5.2), Inches(2.4),
         [[("1.  Checkout reads stock from the database (not the browser)", {"size": 11.5})],
          [("2.  Order placed in one atomic update: ", {"size": 11.5}),
           ("""{"$inc": {"stock": -qty}}""", {"size": 11, "font": "Consolas", "color": SUCCESS})],
          [("3.  If any item is short, the whole order is rejected — ", {"size": 11.5}),
           ("no partial decrements", {"bold": True, "size": 11.5})],
          [("4.  Client cart is cleared only on success", {"size": 11.5})]])
add_rect(s, Inches(0.7), Inches(5.35), Inches(11.93), Inches(1.45), fill=RGBColor(0xEA, 0xF2, 0xFB))
add_text(s, Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.35), "Why this matters", size=12.5, color=ACCENT_DARK, bold=True)
add_text(s, Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.7),
         [[("Two customers hitting the last unit at the same moment: MongoDB's atomic update "
            "means exactly one succeeds. The loser gets a clear \"insufficient stock\" error "
            "with the item name — no negative inventory, no double sales.", {"size": 12})]])
slide_footer(s, 11, TOTAL)

# ---- 12. Orders & tracking ------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "Feature — Order Tracking & Cancellation")
steps = ["pending", "confirmed", "shipped", "delivered"]
x = Inches(0.7)
for i, st in enumerate(steps):
    add_rect(s, x, Inches(1.95), Inches(2.72), Inches(1.35), fill=LIGHT_BG, line=BORDER)
    add_rect(s, x, Inches(1.95), Inches(2.72), Inches(0.12), fill=SUCCESS)
    add_text(s, x + Inches(0.2), Inches(2.2), Inches(2.3), Inches(0.5), st, size=14, color=INK, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.72), Inches(2.3), Inches(0.5),
             "set automatically" if i == 0 else "set by admin",
             size=9.5, color=SUBTLE)
    if i < 3:
        add_rect(s, x + Inches(2.72), Inches(2.35), Inches(0.38), Inches(0.5), fill=SUCCESS, shape=MSO_SHAPE.CHEVRON)
    x += Inches(3.1)
bullets(s, Inches(0.7), Inches(3.7), Inches(6.1), Inches(2.7), [
    ("My Orders page", "every placed order with a live progress tracker"),
    ("Full history", "date, status badge, items, totals — newest first"),
    ("Client cancellation", "allowed while pending or confirmed"),
    ("Stock integrity", "cancelling restores inventory atomically"),
], size=12, gap=9)
add_rect(s, Inches(6.95), Inches(3.7), Inches(5.68), Inches(2.7), fill=INK)
add_text(s, Inches(7.25), Inches(3.9), Inches(5.1), Inches(0.35), "Cancellation flow", size=13, color=CARD, bold=True)
add_text(s, Inches(7.25), Inches(4.35), Inches(5.1), Inches(1.9),
         [[("1.  Customer requests cancel (pending / confirmed only)", {"size": 11.5, "color": MUTED})],
          [("2.  API removes the order", {"size": 11.5, "color": MUTED})],
          [("3.  Stock is added back in the same transaction", {"size": 11.5, "color": MUTED})],
          [("4.  UI shows a success message immediately", {"size": 11.5, "color": MUTED})]])
add_text(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.4),
         [[("Design note: ", {"bold": True, "color": INK}),
           ("the four stages are exactly the values the admin workflow writes — one source of "
            "truth for the whole order lifecycle.", {"color": SUBTLE})]], size=11.5)
slide_footer(s, 12, TOTAL)

# ---- 13. Admin dashboard --------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "Feature — Admin Dashboard & Analytics")
add_rect(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(4.35), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(1.0), Inches(2.05), Inches(5.3), Inches(0.35), "Operations", size=13, color=ACCENT_DARK, bold=True)
bullets(s, Inches(1.0), Inches(2.5), Inches(5.3), Inches(3.6), [
    ("KPI stat cards", "revenue, order count, stock-outs, low stock — from one /stats call"),
    ("Product management", "create, edit, delete, and update stock in place"),
    ("Order workflow", "move any order pending → confirmed → shipped → delivered"),
    ("Search & filter", "find products and orders instantly"),
    ("Role-gated", "route + API both reject non-admins"),
], size=11.5, gap=7)
add_rect(s, Inches(6.85), Inches(1.85), Inches(5.78), Inches(4.35), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(7.15), Inches(2.05), Inches(5.2), Inches(0.35), "Analytics — hand-built SVG charts", size=13, color=ACCENT_DARK, bold=True)
bullets(s, Inches(7.15), Inches(2.5), Inches(5.2), Inches(3.6), [
    ("Revenue last 14 days", "bar chart with value labels & gridlines"),
    ("Orders by status", "color-coded bars mirroring the workflow colors"),
    ("Products by category", "donut with legend, counts and percentages"),
    ("Top products", "horizontal bars by quantity sold + revenue"),
    ("No chart library", "pure SVG components — precise control, zero dependencies"),
], size=11.5, gap=7)
add_text(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
         [[("API:  ", {"bold": True, "color": INK}),
           ("GET /api/admin/stats aggregates orders and products server-side in one round trip.",
            {"color": SUBTLE})]], size=11.5)
slide_footer(s, 13, TOTAL)

# ---- 14. UX & design ------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 2 · BUILD", "UX & Visual Design")
left = [
    ("Light & dark themes", "one toggle in the navbar, saved in localStorage"),
    ("Follows the OS by default", "prefers-color-scheme on first visit"),
    ("No flash on load", "inline script sets the theme before React mounts"),
    ("Design tokens", "196 colors → 38 semantic CSS variables"),
]
add_rect(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(3.3), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(0.35), "Theming", size=13, color=ACCENT_DARK, bold=True)
yy = Inches(2.55)
for head, sub in left:
    add_text(s, Inches(1.0), yy, Inches(5.3), Inches(0.3), f"▸ {head}", size=12, color=INK, bold=True)
    add_text(s, Inches(1.25), yy + Inches(0.26), Inches(5.1), Inches(0.55), sub, size=11, color=SUBTLE)
    yy += Inches(0.72)
add_rect(s, Inches(6.85), Inches(1.9), Inches(5.78), Inches(3.3), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(7.15), Inches(2.1), Inches(5.2), Inches(0.35), "Craft details", size=13, color=ACCENT_DARK, bold=True)
bullets(s, Inches(7.15), Inches(2.55), Inches(5.2), Inches(2.6), [
    ("Responsive layout", "desktop grid, mobile hamburger menu, fluid type"),
    ("Accessibility", "aria labels, keyboard focus rings, semantic buttons"),
    ("State coverage", "loading, empty, error, out-of-stock, success — all styled"),
    ("Consistent feedback", "one error shape across every API failure"),
], size=11.5, gap=8)
add_rect(s, Inches(0.7), Inches(5.45), Inches(11.93), Inches(1.3), fill=INK)
add_text(s, Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.35), "The full palette is token-driven", size=12.5, color=CARD, bold=True)
add_text(s, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.6),
         [[("Themes swap only the token values (", {"size": 11.5, "color": MUTED}),
           ("--bg, --surface, --text…", {"size": 11, "font": "Consolas", "color": ACCENT}),
           (") — components never change, and dark-mode contrast stays ≥ 4.5:1 for body text.",
            {"size": 11.5, "color": MUTED})]])
slide_footer(s, 14, TOTAL)

# ---- 15. Deployment -------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 3 · DELIVERY", "Deployment & DevOps")
add_text(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.9),
         [[("The whole stack runs on one Vercel project: the React build is served as static "
            "files, and the same vercel.json routes /api/* to a Python 3.12 serverless "
            "function. MongoDB runs on Atlas — zero servers to manage.", {"size": 13.5})]],
         line_spacing=1.25)
cols = [
    ("Source control — GitHub", ACCENT_DARK, [
        "private repo, main branch",
        "clear commit history",
        "identity configured for the owner",
    ]),
    ("Hosting — Vercel", ACCENT, [
        "static build (Vite → dist)",
        "Python functions via vercel.json",
        "env vars: MONGODB_URI, JWT_SECRET",
        "auto-deploy on every git push",
    ]),
    ("Database — MongoDB Atlas", SUCCESS, [
        "free cloud cluster",
        "network access for the server",
        "seeded once via script",
    ]),
]
x = Inches(0.7)
for title, accent, rows in cols:
    add_rect(s, x, Inches(2.95), Inches(3.83), Inches(2.9), fill=LIGHT_BG, line=BORDER)
    add_rect(s, x, Inches(2.95), Inches(3.83), Inches(0.55), fill=accent)
    add_text(s, x + Inches(0.25), Inches(3.04), Inches(3.3), Inches(0.4), title, size=13, color=CARD, bold=True)
    yy = Inches(3.7)
    for row in rows:
        add_text(s, x + Inches(0.25), yy, Inches(3.35), Inches(0.4), f"▸ {row}", size=10.5, color=INK)
        yy += Inches(0.45)
    x += Inches(4.0)
add_rect(s, Inches(0.7), Inches(6.05), Inches(11.93), Inches(0.85), fill=RGBColor(0xEA, 0xF2, 0xFB))
add_text(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.5),
         [[("Deploy pipeline:  ", {"bold": True, "color": INK}),
           ("git commit → push to GitHub main → Vercel builds the site + bundles the API → "
            "production alias updates in ~1 minute.", {"color": INK_SOFT})]], size=12)
slide_footer(s, 15, TOTAL)

# ---- 16. Challenges -------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 3 · DELIVERY", "Challenges & How They Were Solved")
problems = [
    ("Products disappeared under category filters",
     "Three seeded products had no category field, so their category was empty and the filter matched nothing.",
     "Added the missing categories in the seed data and re-seeded the database.",
     ACCENT),
    ("Product images broken in production",
     "22 of 25 images pointed at /src/assets paths that only exist on the dev machine, not on Vercel.",
     "Mounted the assets directory in FastAPI and routed /src/assets/* to the backend in vercel.json.",
     SUCCESS),
    ("Vercel blocked every new deployment",
     "Hobby plan rejects deployments whose git commit author isn't the account owner; CLI deploys used a placeholder identity.",
     "Rewrote the commits under the owner's real identity and connected the GitHub repo for auto-deploys.",
     WARNING),
]
y = Inches(1.85)
for i, (head, prob, sol, color) in enumerate(problems):
    add_rect(s, Inches(0.7), y, Inches(0.14), Inches(1.32), fill=color)
    add_text(s, Inches(1.05), y + Inches(0.02), Inches(11.4), Inches(0.35), head, size=13.5, color=INK, bold=True)
    add_text(s, Inches(1.05), y + Inches(0.42), Inches(11.4), Inches(0.5),
             [[("Problem:  ", {"bold": True, "color": DANGER, "size": 11}), (prob, {"size": 11, "color": SUBTLE})]])
    add_text(s, Inches(1.05), y + Inches(0.88), Inches(11.4), Inches(0.5),
             [[("Fix:  ", {"bold": True, "color": SUCCESS, "size": 11}), (sol, {"size": 11, "color": SUBTLE})]])
    y += Inches(1.5)
add_text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.5),
         [[("Lesson: ", {"bold": True, "color": INK}),
           ("data hygiene, environment parity, and identity configuration are the three "
            "silent killers of a demo day.", {"color": SUBTLE})]], size=12)
slide_footer(s, 16, TOTAL)

# ---- 17. Testing ----------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 3 · DELIVERY", "Testing & Verification")
tests = [
    ("Build pipeline", ["npm run build passes cleanly (Vite production build)", "eslint runs with only pre-existing warnings", "Vercel build + Python bundle succeed"], ACCENT),
    ("API smoke tests", ["/api/health, /api/products, /api/auth/login verified on production", "25 products, 0 missing categories after re-seed", "image endpoints return correct MIME types"], SUCCESS),
    ("Manual QA flows", ["register → login → browse → add to cart → checkout", "stock decreases; insufficient-stock rejection shows error", "cancel order → stock restored → tracker disappears"], ACCENT_DARK),
    ("Edge cases covered", ["duplicate signup email → friendly 400", "out-of-stock products can't be added", "admin pages blocked for non-admins", "dark/light theme persists across reloads"], WARNING),
]
x = Inches(0.7)
for title, rows, color in tests:
    add_rect(s, x, Inches(1.9), Inches(5.85), Inches(4.4), fill=LIGHT_BG, line=BORDER)
    add_rect(s, x, Inches(1.9), Inches(5.85), Inches(0.55), fill=color)
    add_text(s, x + Inches(0.25), Inches(1.99), Inches(5.3), Inches(0.4), title, size=13.5, color=CARD, bold=True)
    yy = Inches(2.65)
    for row in rows:
        add_text(s, x + Inches(0.3), yy, Inches(5.3), Inches(0.5), f"▸ {row}", size=10.5, color=INK)
        yy += Inches(0.55)
    x += Inches(6.08) if x < Inches(6) else 0
    if x > Inches(12.5):
        x = Inches(0.7)
        add_text(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.3), "", size=6)
add_text(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5),
         [[("All checks are repeatable: the same commands (build, lint, API calls) are "
            "documented in the repo and were re-run after every fix.", {"size": 11.5, "color": SUBTLE, "italic": True})]])
slide_footer(s, 17, TOTAL)

# ---- 18. Run locally ------------------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 3 · DELIVERY", "Running the Project Locally")
add_rect(s, Inches(0.7), Inches(1.9), Inches(5.85), Inches(4.4), fill=INK)
add_text(s, Inches(1.0), Inches(2.1), Inches(5.2), Inches(0.35), "Three commands", size=14, color=CARD, bold=True)
code_lines = [
    ("npm install", "install frontend dependencies"),
    ("npm run seed", "load 25 products into MongoDB"),
    ("npm run dev:all", "Vite (port 5173) + FastAPI (port 5000) together"),
]
yy = Inches(2.65)
for code, note in code_lines:
    add_rect(s, Inches(1.0), yy, Inches(5.2), Inches(0.75), fill=RGBColor(0x2B, 0x31, 0x3A))
    add_text(s, Inches(1.2), yy + Inches(0.08), Inches(4.9), Inches(0.35), code, size=13, color=ACCENT, font="Consolas", bold=True)
    add_text(s, Inches(1.2), yy + Inches(0.4), Inches(4.9), Inches(0.3), note, size=9.5, color=MUTED)
    yy += Inches(0.95)
add_text(s, Inches(1.0), Inches(5.6), Inches(5.2), Inches(0.6),
         "A .env file holds MONGODB_URI and JWT_SECRET (never committed).", size=10, color=MUTED)
add_rect(s, Inches(6.85), Inches(1.9), Inches(5.78), Inches(4.4), fill=LIGHT_BG, line=BORDER)
add_text(s, Inches(7.15), Inches(2.1), Inches(5.2), Inches(0.35), "Project structure", size=13, color=ACCENT_DARK, bold=True)
tree = [
    ("src/", "React pages, components, contexts, data layer"),
    ("server/", "FastAPI app: routers, schemas, security, seeding"),
    ("api/index.py", "Vercel entry point — re-exports the FastAPI app"),
    ("vercel.json", "builds & routes: static + Python on one domain"),
    ("requirements.txt", "backend dependencies (both local & Vercel)"),
]
yy = Inches(2.6)
for name, note in tree:
    add_text(s, Inches(7.15), yy, Inches(2.2), Inches(0.3), name, size=11, color=ACCENT_DARK, font="Consolas", bold=True)
    add_text(s, Inches(9.4), yy + Inches(0.03), Inches(3.1), Inches(0.55), note, size=10, color=SUBTLE)
    yy += Inches(0.72)
slide_footer(s, 18, TOTAL)

# ---- 19. Future improvements ---------------------------------------------
s = blank_slide(prs)
slide_header(s, "PART 3 · DELIVERY", "Future Improvements")
futures = [
    ("Payments", "Integrate Stripe or PayPal for real checkout instead of order placement"),
    ("Wishlists & reviews", "per-user favorites and product ratings"),
    ("Email notifications", "order confirmation and status change emails"),
    ("Shipping", "address book, courier selection, live tracking numbers"),
    ("Inventory alerts", "low-stock notifications to the admin"),
    ("Performance & scale", "server-side pagination, image CDN, response caching"),
]
x = Inches(0.7)
for i, (head, sub) in enumerate(futures):
    col = i % 3
    row = i // 3
    bx = x + col * Inches(4.05)
    by = Inches(1.95) + row * Inches(2.15)
    add_rect(s, bx, by, Inches(3.83), Inches(1.95), fill=LIGHT_BG, line=BORDER)
    add_rect(s, bx, by, Inches(0.12), Inches(1.95), fill=ACCENT)
    add_text(s, bx + Inches(0.32), by + Inches(0.2), Inches(3.3), Inches(0.4), head, size=13, color=INK, bold=True)
    add_text(s, bx + Inches(0.32), by + Inches(0.7), Inches(3.3), Inches(1.1), sub, size=10.5, color=SUBTLE)
add_text(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.4),
         [[("The architecture is already shaped for these: every feature is a new route, "
            "a new context, or a new collection — no rewrites.", {"size": 11.5, "color": SUBTLE, "italic": True})]])
slide_footer(s, 19, TOTAL)

# ---- 20. Thanks -----------------------------------------------------------
s = blank_slide(prs, bg=INK)
add_rect(s, Inches(0.85), Inches(3.05), Inches(1.1), Inches(0.07), fill=ACCENT)
add_text(s, Inches(0.85), Inches(3.25), Inches(11.6), Inches(1.2),
         "Thank You", size=56, color=CARD, bold=True)
add_text(s, Inches(0.88), Inches(4.4), Inches(11.6), Inches(0.5),
         "Questions & feedback are welcome", size=18, color=MUTED)
add_rect(s, 0, Inches(5.6), SW, Pt(1), fill=RGBColor(0x2B, 0x31, 0x3A))
add_text(s, Inches(0.88), Inches(6.0), Inches(11.6), Inches(1.0),
         [
             [("Live demo:  ", {"bold": True, "color": CARD}),
              ("https://key-forge-final-project.vercel.app", {"color": ACCENT})],
             [("Source:  ", {"bold": True, "color": CARD}),
              ("https://github.com/kimhongpech-dev/KeyForge-Final-Project", {"color": ACCENT})],
         ], size=13)

prs.save(OUT)
print(f"Saved: {OUT} ({OUT.stat().st_size // 1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
